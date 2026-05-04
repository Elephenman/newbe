#!/usr/bin/env python3
"""
组会汇报 PPT 生成器 — 浙大蓝多篇版 v2.0
基于 @PPT架构师 的视觉蓝图 + 浙大蓝模板，生成 .pptx 文件

双路线策略：
  路线A（模板克隆）：打开模板->按layout复制slide->替换文字->删多余slide
  路线B（纯绘制）：  无模板时，用LAYOUT_COORDS实测坐标从空白创建

修复记录：
  P0-1: 所有create_*_slide函数现在接受coords参数并使用LAYOUT_COORDS
  P0-2: LAYOUT_COORDS坐标全部来自tokens.json实测值
  P0-3: create_cover_slide改为浙大蓝F1样式（白底+校徽+标题+底部栏）
  P1-4: --template参数现在有实际功能（模板克隆路线A）
  P1-5: slides_data.json支持image_path字段，实际嵌入论文图片
  P2-9: 增加validate_slides_data() schema验证
  P2-10: F1 LAYOUT_COORDS包含完整元素（分隔线/文献标签/底部栏等）
  P2-13: pt_to_emu()不再是死代码，被add_textbox_at_coords等广泛使用

用法:
  python build_pptx.py --output "output.pptx" --data slides_data.json
  python build_pptx.py --output "output.pptx" --data slides_data.json --template "浙大蓝-多篇版.pptx" --images-dir images/
"""

import json
import argparse
import os
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import sys

# ==================== Layout->坐标映射 ====================
# 所有坐标来自tokens.json实测数据(pt)，需x12700转EMU
LAYOUT_COORDS = {
    'F1-cover': {
        'emblem_group': {'left': 369.89, 'top': 57.52, 'width': 220.22, 'height': 50.18},
        'title': {'left': 105.43, 'top': 164.11, 'width': 749.15, 'height': 60.59},
        'separator_line': {'left': 89.0, 'top': 234.88, 'width': 765.57, 'height': 1.0},
        'label_1': {'left': 180.61, 'top': 265.8, 'width': 86.25, 'height': 29.08},
        'label_2': {'left': 180.61, 'top': 313.8, 'width': 86.25, 'height': 29.08},
        'label_3': {'left': 180.61, 'top': 361.8, 'width': 86.25, 'height': 29.08},
        'lit_title_1': {'left': 271.36, 'top': 265.8, 'width': 583.21, 'height': 29.08},
        'lit_title_2': {'left': 271.36, 'top': 313.8, 'width': 583.21, 'height': 29.08},
        'lit_title_3': {'left': 271.36, 'top': 361.8, 'width': 583.21, 'height': 29.08},
        'bottom_bar': {'left': 247.0, 'top': 454.94, 'width': 466.0, 'height': 37.94},
        'emblem_small': {'left': 278.47, 'top': 462.57, 'width': 22.68, 'height': 22.68},
        'presenter': {'left': 300.9, 'top': 460.58, 'width': 161.34, 'height': 26.66},
        'vertical_line': {'left': 472.72, 'top': 460.58, 'width': 0.75, 'height': 26.66},
        'date_text': {'left': 522.18, 'top': 460.58, 'width': 207.57, 'height': 26.66},
    },
    'F2-toc': {
        'title_area': {'left': 388.52, 'top': 12.66, 'width': 182.95, 'height': 103.81},
        'toc_title': {'left': 432.34, 'top': 36.69, 'width': 101.38, 'height': 55.74},
        'num_1': {'left': 136.81, 'top': 169.53, 'width': 52.64, 'height': 50.14},
        'num_2': {'left': 136.81, 'top': 244.93, 'width': 52.64, 'height': 50.14},
        'num_3': {'left': 136.81, 'top': 320.34, 'width': 52.64, 'height': 50.14},
        'num_4': {'left': 136.81, 'top': 395.74, 'width': 52.64, 'height': 50.14},
        'lit_1': {'left': 195.31, 'top': 178.84, 'width': 627.89, 'height': 31.5},
        'lit_2': {'left': 195.31, 'top': 254.25, 'width': 627.89, 'height': 31.5},
        'lit_3': {'left': 195.31, 'top': 329.65, 'width': 627.89, 'height': 31.5},
        'lit_4': {'left': 195.31, 'top': 405.06, 'width': 627.89, 'height': 31.5},
    },
    'F3-section-separator': {
        'title_area': {'left': 388.52, 'top': 12.66, 'width': 182.95, 'height': 103.81},
        'big_num': {'left': 428.55, 'top': 10.55, 'width': 102.9, 'height': 87.24},
        'part_label': {'left': 427.3, 'top': 86.78, 'width': 105.4, 'height': 24.23},
        'num_1': {'left': 136.81, 'top': 169.53, 'width': 52.64, 'height': 50.14},
        'num_2': {'left': 136.81, 'top': 244.93, 'width': 52.64, 'height': 50.14},
        'num_3': {'left': 136.81, 'top': 320.34, 'width': 52.64, 'height': 50.14},
        'num_4': {'left': 136.81, 'top': 395.74, 'width': 52.64, 'height': 50.14},
        'lit_1': {'left': 195.31, 'top': 178.84, 'width': 627.89, 'height': 31.5},
        'lit_2': {'left': 195.31, 'top': 254.25, 'width': 627.89, 'height': 31.5},
        'lit_3': {'left': 195.31, 'top': 329.65, 'width': 627.89, 'height': 31.5},
        'lit_4': {'left': 195.31, 'top': 405.06, 'width': 627.89, 'height': 31.5},
    },
    'F4-ending': {
        'emblem_group': {'left': 369.89, 'top': 57.52, 'width': 220.22, 'height': 50.18},
        'title': {'left': 105.43, 'top': 200.0, 'width': 749.15, 'height': 80.0},
        'bottom_bar': {'left': 247.0, 'top': 454.94, 'width': 466.0, 'height': 37.94},
        'emblem_small': {'left': 278.47, 'top': 462.57, 'width': 22.68, 'height': 22.68},
        'presenter': {'left': 300.9, 'top': 460.58, 'width': 161.34, 'height': 26.66},
    },
    'C1-info': {
        'section_num': {'left': 29.74, 'top': 32.35, 'width': 51.27, 'height': 36.35},
        'section_title': {'left': 79.55, 'top': 32.35, 'width': 185.0, 'height': 36.35},
        'title_line': {'left': 239.5, 'top': 49.11, 'width': 615.92, 'height': 1.0},
        'left_image': {'left': 54.77, 'top': 103.21, 'width': 562.26, 'height': 273.08},
        'right_table': {'left': 638.23, 'top': 112.21, 'width': 258.77, 'height': 144.0},
        'keyword_1': {'left': 643.46, 'top': 280.82, 'width': 124.77, 'height': 32.25},
        'keyword_2': {'left': 777.46, 'top': 280.82, 'width': 124.77, 'height': 32.25},
        'keyword_3': {'left': 643.46, 'top': 328.04, 'width': 124.77, 'height': 32.25},
        'keyword_4': {'left': 777.46, 'top': 328.04, 'width': 124.77, 'height': 32.25},
        'bottom_area': {'left': 54.77, 'top': 397.38, 'width': 867.23, 'height': 101.7},
        'bottom_label': {'left': 79.55, 'top': 430.06, 'width': 146.0, 'height': 36.35},
        'bottom_text': {'left': 239.5, 'top': 403.38, 'width': 675.5, 'height': 89.71},
        'page_number': {'left': 706, 'top': 500.5, 'width': 216, 'height': 28.75},
    },
    'C2-image-text': {
        'section_num': {'left': 29.74, 'top': 32.35, 'width': 51.27, 'height': 36.35},
        'section_title': {'left': 79.55, 'top': 32.35, 'width': 185.0, 'height': 36.35},
        'title_line': {'left': 264.87, 'top': 49.11, 'width': 590.55, 'height': 1.0},
        'left_image': {'left': 29.74, 'top': 145.44, 'width': 439.01, 'height': 292.2},
        'right_subtitle_1': {'left': 519.56, 'top': 129.69, 'width': 378.19, 'height': 36.35},
        'right_underline_1': {'left': 519.56, 'top': 167.03, 'width': 378.19, 'height': 1.0},
        'right_text_1': {'left': 519.56, 'top': 177.22, 'width': 378.19, 'height': 105.76},
        'right_subtitle_2': {'left': 519.56, 'top': 311.69, 'width': 378.19, 'height': 36.35},
        'right_underline_2': {'left': 519.56, 'top': 349.03, 'width': 378.19, 'height': 1.0},
        'right_text_2': {'left': 519.56, 'top': 359.22, 'width': 378.19, 'height': 105.76},
        'caption': {'left': 164.41, 'top': 462.61, 'width': 169.67, 'height': 24.23},
        'page_number': {'left': 706, 'top': 500.5, 'width': 216, 'height': 28.75},
    },
    'C3-flowchart': {
        'section_num': {'left': 29.74, 'top': 32.35, 'width': 51.27, 'height': 36.35},
        'section_title': {'left': 79.55, 'top': 32.35, 'width': 185.0, 'height': 36.35},
        'title_line': {'left': 239.5, 'top': 49.11, 'width': 615.92, 'height': 1.0},
        'flow_area': {'left': 66, 'top': 103, 'width': 828, 'height': 350},
        'page_number': {'left': 706, 'top': 500.5, 'width': 216, 'height': 28.75},
    },
    'C4-general': {
        'section_num': {'left': 29.74, 'top': 32.35, 'width': 51.27, 'height': 36.35},
        'section_title': {'left': 79.55, 'top': 32.35, 'width': 185.0, 'height': 36.35},
        'title_line': {'left': 239.5, 'top': 49.11, 'width': 615.92, 'height': 1.0},
        'content_area': {'left': 29.74, 'top': 70.0, 'width': 900.0, 'height': 410.0},
        'page_number': {'left': 706, 'top': 500.5, 'width': 216, 'height': 28.75},
    },
}

# ==================== 浙大蓝配色（来自tokens.json） ====================
COLORS = {
    'primary': RGBColor(0x00, 0x3F, 0x88),
    'accent2': RGBColor(0xED, 0x7D, 0x31),
    'accent5': RGBColor(0x5B, 0x9B, 0xD5),
    'accent4': RGBColor(0xFF, 0xC0, 0x00),
    'accent6': RGBColor(0x70, 0xAD, 0x47),
    'dk2': RGBColor(0x44, 0x54, 0x6A),
    'warning': RGBColor(0xC0, 0x00, 0x00),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'light_bg': RGBColor(0xF2, 0xF2, 0xF2),
    'text_dark': RGBColor(0x00, 0x00, 0x00),
    'placeholder_bg': RGBColor(0xD9, 0xD9, 0xD9),
    'border_gray': RGBColor(0xCB, 0xCB, 0xCB),
}

FONTS = {'zh': '微软雅黑', 'en': 'Calibri'}

SLIDE_WIDTH_PT = 960
SLIDE_HEIGHT_PT = 540


def pt_to_emu(pt_val):
    """pt转EMU (1pt = 12700 EMU) — P2-13 FIX: 已被add_textbox_at_coords等广泛使用"""
    return int(pt_val * 12700)


def get_layout_coords(layout_id):
    """获取layout变体的坐标配置"""
    return LAYOUT_COORDS.get(layout_id, LAYOUT_COORDS['C4-general'])


# ==================== 模板克隆工具函数（路线A） ====================

# 模板slide索引映射：type -> 模板中的slide编号
# 浙大蓝多篇版模板 (24 slides):
#   0:封面  1:目录  2:Part1分隔  3:C1-info  4:C2-image-text
#   5:C3-flowchart  6:C4-general(双图)  7:C4-general(图+结论)
#   8:C4-general(三列)  9:Part2分隔  10:C1-info变体
#   11:C2-image-text变体  12:C3-flowchart变体  13:C4-general变体
#   14:Part3分隔  15-18:第三篇文献  19:Part4分隔  20-22:总结
#   23:结尾致谢
TEMPLATE_SLIDE_MAP = {
    'F1-cover': 0,
    'F2-toc': 1,
    'F3-section-separator': 2,
    'C1-info': 3,
    'C2-image-text': 4,
    'C3-flowchart': 5,
    'C4-general': 6,
    'F4-ending': 23,
    # 变体索引(同layout不同样式)
    'C1-info-v2': 10,
    'C2-image-text-v2': 11,
    'C3-flowchart-v2': 12,
    'C4-general-v2': 13,
}

# 模板中的默认文字 -> 需要替换的占位文字
TEMPLATE_PLACEHOLDER_TEXTS = [
    "组会文献汇报PPT模板-多篇文献版",
    "Safety, activity, and immune of anti-PD-1 antibody in cancer",
    "Erlotinib in Previously Treated Non–Small-Cell Lung Cancer",
    "Mutation and Cancer: Statistical Study of Retinoblastoma",
    "文献阅读总结与思考",
    "汇报人：Hason PPT",
    "汇报时间：20XX/XX/XX",
    "文献一", "文献二", "文献三",
    "目录",
    "PART ONE", "PART TWO", "PART THREE", "PART FOUR",
    "01", "02", "03", "04",
    "文献基本信息",
    "研究背景与意义",
    "实验流程及方法",
    "实验结果及结论",
    "文献总结",
    "阅读反思",
    "后续改进方向",
    "恳请各位老师批评指正！",
    "这里输入研究背景",
    "这里输入研究意义",
    "这里输入文献主要研究内容",
    "这里输入实验流程",
    "用图表的数据简单描述实验结果",
    "根据实验结果可以得到什么结论",
    "输入关键词",
    "输入实验/调研的名称",
    "添加标题内容",
    "研究内容",
    "实验方法",
    "实验目的",
    "完成XX实验，获得XX结果",
    "图1  输入图示标题及标注",
    "图2  输入图示标题及标注",
    "图3  输入图示标题及标注",
    "改进方向一", "改进方向二", "改进方向三",
    "改进方向四", "改进方向五", "改进方向六",
    "6大 方向",
    "创新点",
    "借鉴点",
    "与文献研究相关的图片",
    "实验项目名称",
    "校徽批量替换/删除操作步骤",
]


def clone_slide_from_template(prs_target, prs_template, slide_index):
    """从模板中克隆指定索引的slide到目标Presentation"""
    if slide_index >= len(prs_template.slides):
        return None
    template_slide = prs_template.slides[slide_index]
    blank_layout = prs_target.slide_layouts[6]
    new_slide = prs_target.slides.add_slide(blank_layout)
    for shape in template_slide.shapes:
        el = copy.deepcopy(shape._element)
        new_slide.shapes._spTree.append(el)
    return new_slide


def replace_text_in_slide(slide, replacements):
    """在slide中替换文字 — replacements是{旧文字: 新文字}的字典
    
    策略：按旧文字长度降序排序，确保长匹配优先（避免"研究内容"先于"主要研究内容"被匹配）
    对每个段落，从长到短依次尝试替换，一旦替换成功就跳过后续短文本
    """
    # 按旧文字长度降序排序，确保长匹配优先
    sorted_replacements = sorted(replacements.items(), key=lambda x: len(x[0]), reverse=True)
    
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            para_text = para.text
            if not para_text:
                continue
            # 从长到短尝试匹配
            for old_text, new_text in sorted_replacements:
                if not old_text:
                    continue
                if old_text in para_text:
                    # 需要替换 — 保留第一个run的格式
                    if para.runs:
                        first_run = para.runs[0]
                        replaced = para_text.replace(old_text, new_text)
                        first_run.text = replaced
                        # 清空后续run
                        for run in para.runs[1:]:
                            run.text = ""
                    break  # 一旦替换成功就跳过后续短文本


def set_shape_text_by_name(slide, shape_name, new_text, font_size=None, bold=None, color=None):
    """按shape名找到文本框并设置文字，可选择性覆盖字体大小/粗体/颜色"""
    for shape in slide.shapes:
        if shape.name == shape_name and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                if para.runs:
                    para.runs[0].text = new_text
                    if font_size:
                        para.runs[0].font.size = Pt(font_size)
                    if bold is not None:
                        para.runs[0].font.bold = bold
                    if color:
                        para.runs[0].font.color.rgb = color
                else:
                    para.text = new_text
            return True
    return False


def replace_image_in_slide(slide, image_path, shape_name=None):
    """替换slide中指定图片shape的图片内容"""
    if not image_path or not os.path.exists(image_path):
        return False
    target_shape = None
    if shape_name:
        for shape in slide.shapes:
            if shape.name == shape_name:
                target_shape = shape
                break
    if target_shape is None:
        # 尝试找任意图片shape
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture
                target_shape = shape
                break
    if target_shape is None:
        return False
    if hasattr(target_shape, 'image') or target_shape.shape_type == 13:
        left, top, width, height = target_shape.left, target_shape.top, target_shape.width, target_shape.height
        sp = target_shape._element
        sp.getparent().remove(sp)
        slide.shapes.add_picture(image_path, left, top, width, height)
        return True
    return False


def replace_all_images_in_slide(slide, images_dir, image_paths=None):
    """替换slide中所有图片shape — images_dir下找figure_*.png或用image_paths列表"""
    img_idx = 0
    if image_paths is None:
        image_paths = []
    # 收集images_dir下的图片
    if images_dir and os.path.isdir(images_dir):
        import glob
        dir_images = sorted(glob.glob(os.path.join(images_dir, 'figure_*.png')) +
                           glob.glob(os.path.join(images_dir, 'figure_*.jpg')))
        image_paths = list(image_paths) + dir_images

    for shape in list(slide.shapes):
        if shape.shape_type == 13:  # Picture
            if img_idx < len(image_paths) and os.path.exists(image_paths[img_idx]):
                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                sp = shape._element
                sp.getparent().remove(sp)
                slide.shapes.add_picture(image_paths[img_idx], left, top, width, height)
                img_idx += 1


def clear_placeholder_texts(slide):
    """清除模板中的占位提示文字(如'这里输入...'等)和操作提示"""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        full_text = shape.text_frame.text.strip()
        # 清除占位提示文字
        if (full_text.startswith("这里输入") or full_text.startswith("输入") or
            full_text.startswith("添加") or full_text.startswith("用图表的数据简单描述") or
            full_text.startswith("根据实验结果可以得到什么结论") or
            full_text.startswith("校徽批量替换") or
            full_text.startswith("与文献研究相关的图片") or
            full_text == "实验项目名称" or
            full_text == "实验方法" or full_text == "实验目的" or
            full_text == "研究内容" or full_text == "完成XX实验，获得XX结果" or
            full_text.startswith("改进方向") or
            full_text.startswith("创新点") or full_text.startswith("借鉴点") or
            full_text == "6大 方向" or
            "网盘链接" in full_text or
            (full_text.startswith("图") and "输入图示标题" in full_text)):
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = ""


# ==================== Route A: 模板克隆路线主逻辑 ====================

def build_route_a(template_path, slides_data, meta, images_dir=None):
    """Route A: 从模板克隆slide并替换内容，保留模板所有视觉元素
    
    策略：直接打开模板文件 -> 缓存需要的模板slide -> 清空所有slide -> 按需克隆回来 -> 填充内容
    """
    # 先缓存模板slide（从独立加载的模板读取）
    prs_for_cache = Presentation(template_path)
    template_slides_cache = {}
    for idx in range(len(prs_for_cache.slides)):
        template_slides_cache[idx] = prs_for_cache.slides[idx]

    # 确定每个输出slide需要哪个模板slide
    tmpl_indices = []
    for s in slides_data:
        page_type = s.get('type', 'content')
        layout_id = s.get('layout', 'C4-general')
        tmpl_key = _resolve_template_key(page_type, layout_id)
        tmpl_idx = TEMPLATE_SLIDE_MAP.get(tmpl_key, 6)
        tmpl_indices.append(min(tmpl_idx, len(prs_for_cache.slides) - 1))

    # 以模板为基底打开（用文件路径）
    prs = Presentation(template_path)

    # 删除模板原有所有slide（从后往前删避免索引错乱）
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[-1].get(qn('r:id'))
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[-1]

    # 按slides_data逐个克隆模板slide
    # 先提取全局目录项（从TOC页获取）
    global_toc_items = []
    for s in slides_data:
        if s.get('type') == 'toc':
            global_toc_items = s.get('items', s.get('bullets', []))
            break

    for idx, s in enumerate(slides_data):
        page_type = s.get('type', 'content')
        layout_id = s.get('layout', 'C4-general')
        page_num = idx + 1
        total_pages = len(slides_data)

        # 确定克隆哪个模板slide
        tmpl_key = _resolve_template_key(page_type, layout_id)
        tmpl_idx = TEMPLATE_SLIDE_MAP.get(tmpl_key, 6)
        tmpl_idx = min(tmpl_idx, len(template_slides_cache) - 1)
        tmpl_slide = template_slides_cache.get(tmpl_idx)

        # 克隆模板slide的所有shape到新slide
        blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        new_slide = prs.slides.add_slide(blank_layout)
        if tmpl_slide is not None:
            for shape in tmpl_slide.shapes:
                el = copy.deepcopy(shape._element)
                new_slide.shapes._spTree.append(el)

        # ===== 替换文字和图片 =====
        image_path = s.get('image_path', '')
        if image_path and images_dir and not os.path.isabs(image_path):
            image_path = os.path.join(images_dir, image_path)

        if page_type == 'cover':
            _fill_cover_slide_a(new_slide, s, meta)
        elif page_type == 'toc':
            _fill_toc_slide_a(new_slide, s)
        elif page_type == 'section_separator':
            # 自动注入全局目录项到分隔页
            if not s.get('items') and not s.get('bullets') and global_toc_items:
                s = dict(s)  # 复制避免修改原始数据
                s['items'] = global_toc_items
            _fill_section_separator_a(new_slide, s)
        elif page_type == 'ending':
            _fill_ending_slide_a(new_slide, s, meta)
        else:
            _fill_content_slide_a(new_slide, s, page_num, total_pages, images_dir)

    return prs


def _resolve_template_key(page_type, layout_id):
    """根据page_type和layout_id确定模板slide key
    
    批判页/总结页使用C4-general布局，但选择更适合文字内容的模板slide
    """
    if page_type == 'cover':
        return 'F1-cover'
    elif page_type == 'toc':
        return 'F2-toc'
    elif page_type == 'section_separator':
        return 'F3-section-separator'
    elif page_type == 'ending':
        return 'F4-ending'
    elif page_type == 'critique':
        # 批判页用C4-general-v2（slide 13，图+结论布局 → 更适合文字）
        return 'C4-general-v2'
    elif page_type == 'summary':
        # 总结页也用C4-general-v2
        return 'C4-general-v2'
    elif layout_id in TEMPLATE_SLIDE_MAP:
        return layout_id
    else:
        return 'C4-general'


def _fill_cover_slide_a(slide, s, meta):
    """Route A: 填充封面页"""
    title = s.get('title', '')
    subtitle = s.get('subtitle', '')
    presenter = meta.get('presenter', s.get('presenter', ''))
    date_str = meta.get('date', s.get('date_str', ''))

    replacements = {}
    replacements["组会文献汇报PPT模板-多篇文献版"] = title
    replacements["恳请各位老师批评指正！"] = title  # 结尾页标题同位置的fallback

    # 文献标题 — 优先用lit_title_1/2/3，否则从subtitle按|分割
    lit_title_1 = s.get('lit_title_1', '')
    lit_title_2 = s.get('lit_title_2', '')
    lit_title_3 = s.get('lit_title_3', '')
    
    if lit_title_1 or lit_title_2 or lit_title_3:
        lit_items = [lit_title_1, lit_title_2, lit_title_3]
    else:
        lit_items = subtitle.split('|') if subtitle else [title]

    template_titles = [
        "Safety, activity, and immune of anti-PD-1 antibody in cancer",
        "Erlotinib in Previously Treated Non–Small-Cell Lung Cancer",
        "Mutation and Cancer: Statistical Study of Retinoblastoma",
    ]
    lit_labels = ["文献一", "文献二", "文献三"]
    
    # 单篇文献模式：所有文献位置都填同一篇的关键词
    for i in range(3):
        if i < len(lit_items) and lit_items[i]:
            if i < len(template_titles):
                replacements[template_titles[i]] = lit_items[i].strip()
        else:
            # 超出实际文献数量的位置，清空标签和标题
            if i < len(lit_labels):
                replacements[lit_labels[i]] = ""
            if i < len(template_titles):
                replacements[template_titles[i]] = ""

    # 汇报人和日期
    if presenter:
        replacements["汇报人：Hason PPT"] = f"汇报人：{presenter}"
    if date_str:
        replacements["汇报时间：20XX/XX/XX"] = f"汇报时间：{date_str}"

    replace_text_in_slide(slide, replacements)


def _fill_toc_slide_a(slide, s):
    """Route A: 填充目录页"""
    items = s.get('items', s.get('bullets', []))
    replacements = {}

    template_toc_items = [
        "Safety, activity, and immune of anti-PD-1 antibody in cancer",
        "Erlotinib in Previously Treated Non–Small-Cell Lung Cancer",
        "Mutation and Cancer: Statistical Study of Retinoblastoma",
        "文献阅读总结与思考",
    ]
    for i, item in enumerate(items[:4]):
        if i < len(template_toc_items):
            replacements[template_toc_items[i]] = item

    # 清空多余的目录项
    for i in range(len(items), 4):
        if i < len(template_toc_items):
            replacements[template_toc_items[i]] = ""

    replace_text_in_slide(slide, replacements)
    
    # 如果有第5个目录项，添加到04标签旁的文本框
    if len(items) > 4:
        # 找到"04"标签旁的文本框并填入第5项
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if txt == "04":
                    # 找同一行右侧的文本框
                    target_top = shape.top
                    for s2 in slide.shapes:
                        if (s2.has_text_frame and s2.left > shape.left and 
                            abs(s2.top - target_top) < 50000 and
                            not s2.text_frame.text.strip()):
                            _set_text_in_shape(s2, items[4])
                            break
                    break


def _fill_section_separator_a(slide, s):
    """Route A: 填充章节分隔页"""
    section_num = s.get('section_num', '01')
    section_label = s.get('section_label', '')
    items = s.get('items', s.get('bullets', []))
    title = s.get('title', '')

    replacements = {}

    # 替换PART标签
    part_labels = ["PART ONE", "PART TWO", "PART THREE", "PART FOUR"]
    part_text = section_label if section_label else title
    if not part_text:
        try:
            sn = int(section_num)
            nums = ["ONE", "TWO", "THREE", "FOUR"]
            part_text = f"PART {nums[min(sn-1, 3)]}" if sn <= 4 else f"PART {section_num}"
        except:
            part_text = f"PART {section_num}"
    for pl in part_labels:
        replacements[pl] = part_text

    # 替换大号编号 — 只替换标题区域（top < 100pt ≈ 1270000 EMU）的大号编号
    try:
        sn = int(section_num)
        target_num = f"{sn:02d}"
        for shape in slide.shapes:
            if shape.has_text_frame and shape.top < 1270000:  # 只看标题区域
                for para in shape.text_frame.paragraphs:
                    para_text = para.text.strip()
                    if para_text in ("01", "02", "03", "04"):
                        if para.runs:
                            para.runs[0].text = target_num
                            for run in para.runs[1:]:
                                run.text = ""
                        break
    except (ValueError, TypeError):
        pass

    # 替换目录项文字（分隔页右侧也有目录项）
    # 先收集所有内容区域的文本框，按top排序
    content_textboxes = _collect_content_textboxes(slide, min_top=2000000)
    
    # section_separator页如果没有items，就从replacements清空模板文字
    template_toc_items = [
        "Safety, activity, and immune of anti-PD-1 antibody in cancer",
        "Erlotinib in Previously Treated Non–Small-Cell Lung Cancer",
        "Mutation and Cancer: Statistical Study of Retinoblastoma",
        "文献阅读总结与思考",
    ]
    if items:
        for i, item in enumerate(items[:4]):
            if i < len(template_toc_items):
                replacements[template_toc_items[i]] = item
        for i in range(len(items), 4):
            if i < len(template_toc_items):
                replacements[template_toc_items[i]] = ""
    else:
        # 没有items就清空模板文字
        for tti in template_toc_items:
            replacements[tti] = ""

    # 替换section标题标签（如"研究背景"等小标签）
    section_title_candidates = ["研究背景", "实验流程", "实验结果", "文献总结",
                                "阅读反思", "后续改进方向"]
    for stc in section_title_candidates:
        if title and stc != title:
            replacements[stc] = title
        elif title and stc == title:
            pass  # 已经对了
        else:
            replacements[stc] = ""

    replace_text_in_slide(slide, replacements)


def _fill_ending_slide_a(slide, s, meta):
    """Route A: 填充结尾页"""
    title = s.get('title', '恳请各位老师批评指正！')
    presenter = meta.get('presenter', s.get('presenter', ''))
    date_str = meta.get('date', s.get('date_str', ''))

    replacements = {}
    replacements["恳请各位老师批评指正！"] = title
    if presenter:
        replacements["汇报人：Hason PPT"] = f"汇报人：{presenter}"
    if date_str:
        replacements["汇报时间：20XX/XX/XX"] = f"汇报时间：{date_str}"

    replace_text_in_slide(slide, replacements)


def _fill_content_slide_a(slide, s, page_num, total_pages, images_dir=None):
    """Route A: 填充内容页(通用) — 支持所有layout和特殊type"""
    layout_id = s.get('layout', 'C4-general')
    section_num = s.get('section_num', '')
    section_title = s.get('title', '')
    image_path = s.get('image_path', '')
    page_type = s.get('type', 'content')

    # ===== 1. 替换标题栏 =====
    replacements = {}
    # 先替换长文本（避免部分匹配），再替换短文本
    # 按长度降序排列，确保"主要研究内容"在"研究内容"之前被匹配
    title_candidates = [
        "主要研究内容",  # 必须在"研究内容"之前
        "文献基本信息", "研究背景与意义", "实验流程及方法",
        "实验结果及结论", "文献总结", "阅读反思", "后续改进方向",
        "研究背景", "研究意义", "实验方法", "实验目的", "研究内容",
    ]
    for tc in title_candidates:
        replacements[tc] = section_title

    # 替换section编号
    num_candidates = ["1.1", "1.2", "1.3", "1.4", "2.1", "2.2", "2.3", "2.4",
                      "3.1", "3.2", "3.3", "3.4", "4.1", "4.2", "4.3"]
    for nc in num_candidates:
        replacements[nc] = section_num

    # 替换副标题占位
    if layout_id == 'C2-image-text':
        bullets = s.get('bullets', [])
        if len(bullets) > 0:
            # 根据section_title推断副标题
            sec_lower = section_title.lower() if section_title else ''
            if '鉴定' in section_title or '扫描' in section_title or 'motif紧邻' in section_title:
                sub1, sub2 = "核心发现", "转录验证"
            elif '突变' in section_title or '启动子' in section_title or '替换' in section_title:
                sub1, sub2 = "实验验证", "功能分析"
            elif '结合' in section_title or 'RpoD' in section_title or '相互作用' in section_title:
                sub1, sub2 = "分子互作", "结构基础"
            elif '保守' in section_title or 'D-T门' in section_title or '广泛' in section_title:
                sub1, sub2 = "跨物种分析", "进化保守性"
            else:
                sub1, sub2 = "研究发现", "实验验证"
            replacements["研究背景"] = sub1
            replacements["研究意义"] = sub2
        # 替换图片caption
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if "输入图示标题" in txt or (txt.startswith("图") and len(txt) < 30):
                    for para in shape.text_frame.paragraphs:
                        if para.runs:
                            para.runs[0].text = f"图  {section_title[:20]}"
                            for r in para.runs[1:]:
                                r.text = ""

    replace_text_in_slide(slide, replacements)

    # ===== 2. 清除占位文字 =====
    clear_placeholder_texts(slide)

    # ===== 3. 替换图片 =====
    abs_image_path = ''
    if image_path:
        abs_image_path = image_path
        if images_dir and not os.path.isabs(image_path):
            abs_image_path = os.path.join(images_dir, image_path)

    if abs_image_path and os.path.exists(abs_image_path):
        replaced = False
        for shape in list(slide.shapes):
            if shape.shape_type == 13:
                if not replaced:
                    left, top, width, height = shape.left, shape.top, shape.width, shape.height
                    sp = shape._element
                    sp.getparent().remove(sp)
                    slide.shapes.add_picture(abs_image_path, left, top, width, height)
                    replaced = True

    # ===== 4. 填充内容文字 =====
    bullets = s.get('bullets', [])
    content_text = s.get('content_text', '')
    if not content_text and bullets:
        content_text = '\n'.join(f"- {b}" for b in bullets)

    if content_text:
        _fill_content_text(slide, content_text, layout_id)

    # ===== 5. 特殊type处理 =====
    if page_type == 'critique':
        _fill_critique_slide_a(slide, s)
    elif page_type == 'summary':
        _fill_summary_slide_a(slide, s)

    # ===== 6. 布局特有填充 =====
    if layout_id == 'C1-info':
        _fill_info_specific(slide, s)
    elif layout_id == 'C3-flowchart':
        _fill_flowchart_specific(slide, s)

    # ===== 7. 更新页码 =====
    _update_page_number(slide, page_num, total_pages)


def _fill_content_text(slide, content_text, layout_id):
    """将内容文字填入slide中合适的空白文本框
    
    策略：收集所有文本框，按top排序，找到内容区域（标题栏以下）的空文本框，
    将内容分发到这些文本框中。不再限制只找"文本框"开头的shape。
    """
    # 收集所有文本框（排除标题栏和页码）
    content_textboxes = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        # 跳过标题栏区域的shape（top < 70pt ≈ 889000 EMU）
        if shape.top < 889000:
            continue
        # 跳过页码区域
        if text.isdigit() and int(text) < 100 and shape.top > 6000000:
            continue
        # 收集内容区域的文本框（空或占位文字）
        if (not text or text.startswith("这里输入") or text.startswith("输入") or
            text.startswith("添加") or text.startswith("用图表的数据简单描述") or
            text.startswith("根据实验结果可以得到什么结论") or
            text.startswith("研究背景") or text.startswith("研究意义") or
            text.startswith("实验方法") or text.startswith("实验目的") or
            text.startswith("研究内容") or text.startswith("完成XX") or
            text.startswith("图") and "输入图示标题" in text):
            content_textboxes.append(shape)

    # 按top位置排序
    content_textboxes.sort(key=lambda s: s.top)
    
    if not content_textboxes:
        return

    # 将bullets分行分发
    lines = content_text.split('\n')
    
    if layout_id == 'C2-image-text' and len(content_textboxes) >= 2:
        # C2布局：右侧有2个文字区域，把bullets分2组
        mid = len(lines) // 2
        if len(lines) <= 2:
            # bullets太少，全放第一个
            _set_text_in_shape(content_textboxes[0], '\n'.join(lines))
        else:
            _set_text_in_shape(content_textboxes[0], '\n'.join(lines[:mid]))
            if len(content_textboxes) > 1:
                _set_text_in_shape(content_textboxes[1], '\n'.join(lines[mid:]))
    elif layout_id == 'C3-flowchart' and content_textboxes:
        # C3布局：流程图下方有描述文字区域
        _set_text_in_shape(content_textboxes[0], '\n'.join(lines))
    elif layout_id == 'C1-info' and len(content_textboxes) >= 1:
        # C1布局：底部描述区
        _set_text_in_shape(content_textboxes[0], '\n'.join(lines))
    else:
        # 通用布局：所有内容放第一个可用文本框
        _set_text_in_shape(content_textboxes[0], '\n'.join(lines))


def _set_text_in_shape(shape, text):
    """安全地设置shape的文本，保留格式"""
    if not shape.has_text_frame:
        return
    # 用第一个paragraph设置文本
    paras = shape.text_frame.paragraphs
    if paras:
        runs = paras[0].runs
        if runs:
            runs[0].text = text
            # 清空后续run
            for r in runs[1:]:
                r.text = ""
        else:
            paras[0].text = text


def _update_page_number(slide, page_num, total_pages):
    """更新页码 — 找底部小数字文本框替换"""
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            # 页码特征：底部、小数字
            is_page_num = False
            if text.isdigit() and int(text) < 100:
                is_page_num = True
            elif '/' in text and any(c.isdigit() for c in text):
                # 格式如 "4 / 10"
                is_page_num = True
            
            if is_page_num and shape.top > 5000000:  # 下半部分
                page_text = f"{page_num} / {total_pages}"
                for para in shape.text_frame.paragraphs:
                    if para.runs:
                        para.runs[0].text = page_text
                        for run in para.runs[1:]:
                            run.text = ""
                    else:
                        para.text = page_text
                return


def _fill_info_specific(slide, s):
    """C1-info特有：关键词、底部文字"""
    keywords = s.get('keywords', [])
    
    # 替换关键词shape
    kw_shapes = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text == "输入关键词" or text.startswith("关键词"):
                kw_shapes.append(shape)
    
    # 也找圆角矩形里的关键词位
    if not kw_shapes:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.top > 3500000:  # 下半部分
                text = shape.text_frame.text.strip()
                if not text or text.startswith("输入"):
                    kw_shapes.append(shape)

    for i, kw in enumerate(keywords[:len(kw_shapes)]):
        _set_text_in_shape(kw_shapes[i], kw)


def _fill_image_text_specific(slide, s):
    """C2-image-text特有：右侧文字 — 已合并到_fill_content_slide_a中处理"""
    pass


def _fill_critique_slide_a(slide, s):
    """Route A: 填充批判页 — 找到内容区域文本框填入批判内容"""
    critiques = s.get('critiques', [])
    if not critiques:
        return
    
    # 收集内容区域空文本框
    content_boxes = _collect_content_textboxes(slide)
    
    # 构建批判内容文本
    lines = []
    for i, c in enumerate(critiques):
        sev = c.get('severity', 'P1')
        title = c.get('title', '')
        detail = c.get('detail', '')
        lines.append(f"[{sev}] {title}")
        if detail:
            lines.append(f"  → {detail}")
    
    critique_text = '\n'.join(lines)
    
    if content_boxes:
        _set_text_in_shape(content_boxes[0], critique_text)


def _fill_summary_slide_a(slide, s):
    """Route A: 填充总结页 — 找到内容区域文本框填入总结内容"""
    conclusions = s.get('conclusions', [])
    inspirations = s.get('inspirations', [])
    
    # 收集内容区域空文本框
    content_boxes = _collect_content_textboxes(slide)
    
    # 构建总结文本
    lines = []
    if conclusions:
        lines.append("核心结论：")
        for c in conclusions:
            lines.append(f"  • {c}")
    if inspirations:
        lines.append("")
        lines.append("对本课题的启发：")
        for ins in inspirations:
            lines.append(f"  • {ins}")
    
    summary_text = '\n'.join(lines)
    
    if content_boxes:
        _set_text_in_shape(content_boxes[0], summary_text)


def _collect_content_textboxes(slide, min_top=889000):
    """收集slide中内容区域的空/可覆写文本框（排除标题栏和页码）
    
    Args:
        slide: pptx slide对象
        min_top: 最小top位置(EMU)，默认889000=70pt，排除标题栏
    
    Returns:
        按top排序的shape列表
    """
    boxes = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.top < min_top:
            continue
        text = shape.text_frame.text.strip()
        # 跳过页码
        if text.isdigit() and int(text) < 100 and shape.top > 6000000:
            continue
        # 收集空文本框或占位文字
        if (not text or text.startswith("这里输入") or text.startswith("输入") or
            text.startswith("添加") or text.startswith("用图表") or
            text.startswith("根据实验") or text.startswith("研究背景") or
            text.startswith("研究意义") or text.startswith("实验方法") or
            text.startswith("实验目的") or text.startswith("研究内容") or
            text.startswith("完成XX") or
            (text.startswith("图") and "输入图示标题" in text) or
            text.startswith("创新点") or text.startswith("借鉴点") or
            text.startswith("改进方向")):
            boxes.append(shape)
    
    boxes.sort(key=lambda s: s.top)
    return boxes


def _fill_flowchart_specific(slide, s):
    """C3-flowchart特有：替换流程节点文字"""
    flow_items = s.get('flow_items', s.get('bullets', []))
    flowchart_text = s.get('flowchart_text', '')

    # 替换流程描述文字
    if flowchart_text:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text.startswith("这里输入实验流程") or text.startswith("这里输入文献"):
                    _set_text_in_shape(shape, flowchart_text)
                    break

    # 替换流程节点（椭圆形状中的文字如"实验 方法"）
    # 按left位置排序找到流程节点
    node_shapes = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip().replace('\n', '').replace(' ', '')
            # 流程节点的文字通常是"实验方法"/"研究内容"/"实验目的"等占位
            if text in ("实验方法", "研究内容", "实验目的", 
                       "添加标题内容", "完成XX实验，获得XX结果"):
                node_shapes.append(shape)
    
    # 也找TextBox类的节点（Auto_shape with text like "研究内容"）
    for shape in slide.shapes:
        if shape.has_text_frame and shape.shape_type == 1:  # AUTO_SHAPE
            text = shape.text_frame.text.strip().replace('\n', '').replace(' ', '')
            if text in ("研究内容",) and shape not in node_shapes:
                node_shapes.append(shape)
    
    # 按left排序
    node_shapes.sort(key=lambda s: s.left)
    
    # 如果有bullets，替换节点文字
    if flow_items and node_shapes:
        for i, item in enumerate(flow_items[:len(node_shapes)]):
            _set_text_in_shape(node_shapes[i], item)


# ==================== 纯绘制工具函数（路线B） ====================

def add_shape_at_coords(slide, shape_type, coords_key, coords, **kwargs):
    c = coords.get(coords_key, {})
    left = pt_to_emu(c.get('left', 30))
    top = pt_to_emu(c.get('top', 30))
    width = pt_to_emu(c.get('width', 800))
    height = pt_to_emu(c.get('height', 50))
    return slide.shapes.add_shape(shape_type, left, top, width, height)


def add_textbox_at_coords(slide, coords_key, coords, text, font_size=16, bold=False,
                          color=None, alignment=PP_ALIGN.LEFT, font_name=None):
    c = coords.get(coords_key, {})
    left = pt_to_emu(c.get('left', 30))
    top = pt_to_emu(c.get('top', 30))
    width = pt_to_emu(c.get('width', 800))
    height = pt_to_emu(c.get('height', 50))
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.name = font_name or FONTS['zh']
    if color:
        p.font.color.rgb = color
    p.alignment = alignment
    return tb


def add_image_at_coords(slide, coords_key, coords, image_path, placeholder_desc="插入图片"):
    """P1-5 FIX: 有图则嵌入，无图则占位"""
    c = coords.get(coords_key, {})
    left = pt_to_emu(c.get('left', 30))
    top = pt_to_emu(c.get('top', 103))
    width = pt_to_emu(c.get('width', 500))
    height = pt_to_emu(c.get('height', 300))
    if image_path and os.path.exists(image_path):
        slide.shapes.add_picture(image_path, left, top, width, height)
    else:
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = COLORS['placeholder_bg']
        box.line.color.rgb = COLORS['border_gray']
        box.line.width = Pt(1)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = placeholder_desc
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
        p.alignment = PP_ALIGN.CENTER


def add_line_at_coords(slide, coords_key, coords, color=None):
    c = coords.get(coords_key, {})
    left = pt_to_emu(c.get('left', 30))
    top = pt_to_emu(c.get('top', 49))
    width = pt_to_emu(c.get('width', 800))
    height = max(pt_to_emu(c.get('height', 1)), pt_to_emu(1))
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    line.fill.solid()
    line.fill.fore_color.rgb = color or COLORS['primary']
    line.line.fill.background()


def add_page_number(slide, coords, page_num, total=10):
    add_textbox_at_coords(slide, 'page_number', coords, f"{page_num} / {total}",
                          font_size=14, color=COLORS['dk2'], alignment=PP_ALIGN.RIGHT)


def add_title_bar_zju(slide, section_num, section_title, coords):
    """添加浙大蓝标准标题栏"""
    add_textbox_at_coords(slide, 'section_num', coords, section_num,
                          font_size=24, bold=True, color=COLORS['primary'])
    add_textbox_at_coords(slide, 'section_title', coords, section_title,
                          font_size=24, bold=True, color=COLORS['primary'])
    add_line_at_coords(slide, 'title_line', coords, COLORS['primary'])


# ==================== 各Layout变体创建函数 ====================

def create_cover_slide(prs, title, subtitle="", presenter="", date_str="", coords=None):
    """F1-cover: 浙大蓝封面页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if coords is None:
        coords = get_layout_coords('F1-cover')

    # 校徽占位
    add_shape_at_coords(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 'emblem_group', coords)
    c = coords.get('emblem_group', {})
    add_textbox_at_coords(slide, 'emblem_group', coords, "浙江大学",
                          font_size=18, bold=True, color=COLORS['primary'], alignment=PP_ALIGN.CENTER)

    # 大标题
    add_textbox_at_coords(slide, 'title', coords, title,
                          font_size=40, bold=True, color=COLORS['primary'], alignment=PP_ALIGN.LEFT)

    # 分隔线
    add_line_at_coords(slide, 'separator_line', coords, COLORS['primary'])

    # 文献标签+标题
    lit_items = subtitle.split('|') if subtitle else [title]
    for i, item in enumerate(lit_items[:3], 1):
        tag_key = f'label_{i}'
        title_key = f'lit_title_{i}'
        if tag_key in coords:
            tag = add_shape_at_coords(slide, MSO_SHAPE.ROUNDED_RECTANGLE, tag_key, coords)
            tag.fill.solid()
            tag.fill.fore_color.rgb = COLORS['primary']
            tag.line.fill.background()
            tf = tag.text_frame
            p = tf.paragraphs[0]
            p.text = f"文献{'一二三'[i-1]}"
            p.font.size = Pt(12)
            p.font.color.rgb = COLORS['white']
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
        if title_key in coords:
            add_textbox_at_coords(slide, title_key, coords, item.strip(),
                                  font_size=14, color=COLORS['dk2'])

    # 底部栏
    bar = add_shape_at_coords(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 'bottom_bar', coords)
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS['primary']
    bar.line.fill.background()

    # 底部小校徽占位
    add_shape_at_coords(slide, MSO_SHAPE.OVAL, 'emblem_small', coords)

    if presenter:
        add_textbox_at_coords(slide, 'presenter', coords, presenter,
                              font_size=14, color=COLORS['white'])
    if date_str:
        add_textbox_at_coords(slide, 'date_text', coords, date_str,
                              font_size=14, color=COLORS['white'])
    return slide


def create_toc_slide(prs, items, coords=None):
    """F2-toc: 目录页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if coords is None:
        coords = get_layout_coords('F2-toc')
    area = add_shape_at_coords(slide, MSO_SHAPE.RECTANGLE, 'title_area', coords)
    area.fill.solid()
    area.fill.fore_color.rgb = COLORS['primary']
    area.line.fill.background()
    add_textbox_at_coords(slide, 'toc_title', coords, '目录',
                          font_size=40, bold=True, color=COLORS['white'], alignment=PP_ALIGN.CENTER)
    for i, item in enumerate(items[:4], 1):
        num_key, lit_key = f'num_{i}', f'lit_{i}'
        if num_key in coords:
            num_box = add_shape_at_coords(slide, MSO_SHAPE.ROUNDED_RECTANGLE, num_key, coords)
            num_box.fill.solid()
            num_box.fill.fore_color.rgb = COLORS['primary']
            num_box.line.fill.background()
            ntf = num_box.text_frame
            np_ = ntf.paragraphs[0]
            np_.text = f"0{i}"
            np_.font.size = Pt(28)
            np_.font.bold = True
            np_.font.color.rgb = COLORS['white']
            np_.alignment = PP_ALIGN.CENTER
        if lit_key in coords:
            add_textbox_at_coords(slide, lit_key, coords, item,
                                  font_size=18, bold=True, color=COLORS['dk2'])
    return slide


def create_section_separator_slide(prs, section_num, section_label, items, coords=None):
    """F3-section-separator: 章节分隔页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if coords is None:
        coords = get_layout_coords('F3-section-separator')
    area = add_shape_at_coords(slide, MSO_SHAPE.RECTANGLE, 'title_area', coords)
    area.fill.solid()
    area.fill.fore_color.rgb = COLORS['primary']
    area.line.fill.background()
    add_textbox_at_coords(slide, 'big_num', coords, f"0{section_num}",
                          font_size=66, bold=True, color=COLORS['white'], alignment=PP_ALIGN.CENTER)
    add_textbox_at_coords(slide, 'part_label', coords, section_label,
                          font_size=14, color=COLORS['white'], alignment=PP_ALIGN.CENTER)
    for i, item in enumerate(items[:4], 1):
        num_key, lit_key = f'num_{i}', f'lit_{i}'
        if num_key in coords:
            num_box = add_shape_at_coords(slide, MSO_SHAPE.ROUNDED_RECTANGLE, num_key, coords)
            num_box.fill.solid()
            num_box.fill.fore_color.rgb = COLORS['primary']
            num_box.line.fill.background()
            ntf = num_box.text_frame
            np_ = ntf.paragraphs[0]
            np_.text = f"0{i}"
            np_.font.size = Pt(28)
            np_.font.bold = True
            np_.font.color.rgb = COLORS['white']
            np_.alignment = PP_ALIGN.CENTER
        if lit_key in coords:
            add_textbox_at_coords(slide, lit_key, coords, item,
                                  font_size=18, bold=True, color=COLORS['dk2'])
    return slide


def create_ending_slide(prs, coords=None):
    """F4-ending: 结尾致谢页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if coords is None:
        coords = get_layout_coords('F4-ending')
    add_shape_at_coords(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 'emblem_group', coords)
    add_textbox_at_coords(slide, 'title', coords, '谢谢聆听',
                          font_size=44, bold=True, color=COLORS['primary'], alignment=PP_ALIGN.CENTER)
    bar = add_shape_at_coords(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 'bottom_bar', coords)
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS['primary']
    bar.line.fill.background()
    return slide


def create_info_slide(prs, section_num, section_title, image_path=None,
                      table_data=None, keywords=None, bottom_text="",
                      page_num=1, total=10, coords=None):
    """C1-info: 文献基本信息页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if coords is None:
        coords = get_layout_coords('C1-info')
    add_title_bar_zju(slide, section_num, section_title, coords)
    add_image_at_coords(slide, 'left_image', coords, image_path, "论文Figure")
    add_textbox_at_coords(slide, 'right_table', coords, table_data or "信息表格",
                          font_size=14, color=COLORS['dk2'] if table_data else RGBColor(0x80,0x80,0x80))
    if keywords:
        for i, kw in enumerate(keywords[:4], 1):
            kw_key = f'keyword_{i}'
            if kw_key in coords:
                kw_box = add_shape_at_coords(slide, MSO_SHAPE.ROUNDED_RECTANGLE, kw_key, coords)
                kw_box.fill.solid()
                kw_box.fill.fore_color.rgb = COLORS['accent5']
                kw_box.line.fill.background()
                ktf = kw_box.text_frame
                kp = ktf.paragraphs[0]
                kp.text = kw
                kp.font.size = Pt(12)
                kp.font.color.rgb = COLORS['white']
                kp.alignment = PP_ALIGN.CENTER
    if bottom_text:
        add_textbox_at_coords(slide, 'bottom_text', coords, bottom_text,
                              font_size=16, color=COLORS['dk2'])
    add_page_number(slide, coords, page_num, total)
    return slide


def create_image_text_slide(prs, section_num, section_title, image_path=None,
                            right_col1_title="", right_col1_text="",
                            right_col2_title="", right_col2_text="",
                            caption="", page_num=1, total=10, coords=None):
    """C2-image-text: 左图右文页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if coords is None:
        coords = get_layout_coords('C2-image-text')
    add_title_bar_zju(slide, section_num, section_title, coords)
    add_image_at_coords(slide, 'left_image', coords, image_path, "论文Figure")
    if right_col1_title:
        add_textbox_at_coords(slide, 'right_subtitle_1', coords, right_col1_title,
                              font_size=22, bold=True, color=COLORS['primary'])
        add_line_at_coords(slide, 'right_underline_1', coords, COLORS['primary'])
    if right_col1_text:
        add_textbox_at_coords(slide, 'right_text_1', coords, right_col1_text,
                              font_size=16, color=COLORS['dk2'])
    if right_col2_title:
        add_textbox_at_coords(slide, 'right_subtitle_2', coords, right_col2_title,
                              font_size=22, bold=True, color=COLORS['primary'])
        add_line_at_coords(slide, 'right_underline_2', coords, COLORS['primary'])
    if right_col2_text:
        add_textbox_at_coords(slide, 'right_text_2', coords, right_col2_text,
                              font_size=16, color=COLORS['dk2'])
    if caption:
        add_textbox_at_coords(slide, 'caption', coords, caption,
                              font_size=14, color=COLORS['dk2'], alignment=PP_ALIGN.CENTER)
    add_page_number(slide, coords, page_num, total)
    return slide


def create_flowchart_slide(prs, section_num, section_title, flow_items=None,
                           page_num=1, total=10, coords=None):
    """C3-flowchart: 流程图页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if coords is None:
        coords = get_layout_coords('C3-flowchart')
    add_title_bar_zju(slide, section_num, section_title, coords)
    fc = coords.get('flow_area', {})
    flow_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        pt_to_emu(fc.get('left',66)), pt_to_emu(fc.get('top',103)),
        pt_to_emu(fc.get('width',828)), pt_to_emu(fc.get('height',350)))
    flow_box.fill.solid()
    flow_box.fill.fore_color.rgb = COLORS['light_bg']
    flow_box.line.color.rgb = COLORS['border_gray']
    flow_box.line.width = Pt(1)
    if flow_items:
        n = len(flow_items)
        node_w = min(120, 828 // (n * 2))
        start_x = 66 + (828 - n * node_w * 2) // 2
        y_center = 103 + 175
        for i, item in enumerate(flow_items[:5]):
            x = start_x + i * node_w * 2
            ellipse = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                pt_to_emu(x), pt_to_emu(y_center-54), pt_to_emu(node_w), pt_to_emu(108))
            ellipse.fill.solid()
            ellipse.fill.fore_color.rgb = COLORS['primary']
            ellipse.line.fill.background()
            etf = ellipse.text_frame
            etf.word_wrap = True
            ep = etf.paragraphs[0]
            ep.text = str(item)
            ep.font.size = Pt(12)
            ep.font.color.rgb = COLORS['white']
            ep.alignment = PP_ALIGN.CENTER
            if i < n - 1:
                arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                    pt_to_emu(x+node_w+5), pt_to_emu(y_center-8),
                    pt_to_emu(node_w-10), pt_to_emu(16))
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = COLORS['accent2']
                arrow.line.fill.background()
    add_page_number(slide, coords, page_num, total)
    return slide


def create_general_content_slide(prs, section_num, section_title, content_text="",
                                 image_path=None, page_num=1, total=10, coords=None):
    """C4-general: 通用内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if coords is None:
        coords = get_layout_coords('C4-general')
    add_title_bar_zju(slide, section_num, section_title, coords)
    if content_text:
        add_textbox_at_coords(slide, 'content_area', coords, content_text,
                              font_size=16, color=COLORS['dk2'])
    if image_path:
        add_image_at_coords(slide, 'content_area', coords, image_path, "图表/结果")
    add_page_number(slide, coords, page_num, total)
    return slide


def create_critique_slide(prs, critiques, fix_suggestions="",
                          section_num="", section_title="方法论批判与局限",
                          page_num=1, total=10, coords=None):
    """批判页：红色警示风格"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if coords is None:
        coords = get_layout_coords('C4-general')
    c = coords.get('section_num', {})
    num_tb = slide.shapes.add_textbox(
        pt_to_emu(c.get('left',30)), pt_to_emu(c.get('top',32)),
        pt_to_emu(c.get('width',51)), pt_to_emu(c.get('height',36)))
    np_ = num_tb.text_frame.paragraphs[0]
    np_.text = section_num or f"{page_num}"
    np_.font.size = Pt(24)
    np_.font.bold = True
    np_.font.color.rgb = COLORS['warning']
    c = coords.get('section_title', {})
    title_tb = slide.shapes.add_textbox(
        pt_to_emu(c.get('left',80)), pt_to_emu(c.get('top',32)),
        pt_to_emu(c.get('width',185)), pt_to_emu(c.get('height',36)))
    tp = title_tb.text_frame.paragraphs[0]
    tp.text = f"⚠️ {section_title}"
    tp.font.size = Pt(24)
    tp.font.bold = True
    tp.font.color.rgb = COLORS['warning']
    add_line_at_coords(slide, 'title_line', coords, COLORS['warning'])
    card_top = 103
    for i, critique in enumerate(critiques[:3]):
        card_y = card_top + i * 140
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            pt_to_emu(66), pt_to_emu(card_y), pt_to_emu(828), pt_to_emu(120))
        card.line.color.rgb = COLORS['warning']
        card.line.width = Pt(2)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0xFF, 0xF0, 0xF0)
        crit_tb = slide.shapes.add_textbox(
            pt_to_emu(80), pt_to_emu(card_y+10), pt_to_emu(800), pt_to_emu(100))
        ctf = crit_tb.text_frame
        ctf.word_wrap = True
        cp = ctf.paragraphs[0]
        cp.text = f"🔴 致命伤 #{i+1}：{critique.get('title', '')}"
        cp.font.size = Pt(18)
        cp.font.bold = True
        cp.font.color.rgb = COLORS['warning']
        dp = ctf.add_paragraph()
        dp.text = critique.get('detail', '')
        dp.font.size = Pt(15)
        dp.font.color.rgb = COLORS['dk2']
        dp.space_before = Pt(6)
    if fix_suggestions:
        fix_tb = slide.shapes.add_textbox(
            pt_to_emu(66), pt_to_emu(card_top + len(critiques)*140 + 10),
            pt_to_emu(828), pt_to_emu(40))
        fp = fix_tb.text_frame.paragraphs[0]
        fp.text = f"💊 修补方向：{fix_suggestions}"
        fp.font.size = Pt(14)
        fp.font.color.rgb = COLORS['accent6']
    add_page_number(slide, coords, page_num, total)
    return slide


def create_summary_slide(prs, takeaway, inspirations, next_steps,
                         section_num="", section_title="总结与下一步",
                         page_num=1, total=10, coords=None):
    """总结页：浙大蓝风格"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if coords is None:
        coords = get_layout_coords('C4-general')
    add_title_bar_zju(slide, section_num or f"{page_num}", f"✅ {section_title}", coords)
    tk_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        pt_to_emu(150), pt_to_emu(103), pt_to_emu(660), pt_to_emu(80))
    tk_box.line.color.rgb = COLORS['primary']
    tk_box.line.width = Pt(3)
    tk_box.fill.solid()
    tk_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xFE)
    tk_tb = slide.shapes.add_textbox(pt_to_emu(170), pt_to_emu(115), pt_to_emu(620), pt_to_emu(56))
    ttf = tk_tb.text_frame
    ttf.word_wrap = True
    tp = ttf.paragraphs[0]
    tp.text = f"📌  一句话总结：{takeaway}"
    tp.font.size = Pt(20)
    tp.font.bold = True
    tp.font.color.rgb = COLORS['primary']
    tp.alignment = PP_ALIGN.CENTER
    y = 210
    if inspirations:
        insp_tb = slide.shapes.add_textbox(pt_to_emu(66), pt_to_emu(y), pt_to_emu(420), pt_to_emu(200))
        itf = insp_tb.text_frame
        itf.word_wrap = True
        ip = itf.paragraphs[0]
        ip.text = "→ 对本课题组的启发"
        ip.font.size = Pt(18)
        ip.font.bold = True
        ip.font.color.rgb = COLORS['accent6']
        for insp in inspirations:
            ip2 = itf.add_paragraph()
            ip2.text = f"  ·  {insp}"
            ip2.font.size = Pt(15)
            ip2.font.color.rgb = COLORS['dk2']
            ip2.space_before = Pt(6)
    if next_steps:
        next_tb = slide.shapes.add_textbox(pt_to_emu(500), pt_to_emu(y), pt_to_emu(420), pt_to_emu(200))
        ntf = next_tb.text_frame
        ntf.word_wrap = True
        np_ = ntf.paragraphs[0]
        np_.text = "→ 下一步行动"
        np_.font.size = Pt(18)
        np_.font.bold = True
        np_.font.color.rgb = COLORS['primary']
        for step in next_steps:
            np2 = ntf.add_paragraph()
            np2.text = f"  ·  {step}"
            np2.font.size = Pt(15)
            np2.font.color.rgb = COLORS['dk2']
            np2.space_before = Pt(6)
    thanks_tb = slide.shapes.add_textbox(pt_to_emu(66), pt_to_emu(440), pt_to_emu(828), pt_to_emu(40))
    thp = thanks_tb.text_frame.paragraphs[0]
    thp.text = "✨  感谢聆听，欢迎讨论！"
    thp.font.size = Pt(22)
    thp.font.color.rgb = COLORS['primary']
    thp.alignment = PP_ALIGN.CENTER
    add_page_number(slide, coords, page_num, total)
    return slide


# ==================== JSON Schema验证（P2-9 FIX） ====================

SLIDES_DATA_SCHEMA = {
    'required_per_slide': ['type', 'layout', 'title'],
    'valid_types': ['cover', 'toc', 'section_separator', 'content', 'critique', 'summary', 'ending'],
    'valid_layouts': list(LAYOUT_COORDS.keys()),
}


def validate_slides_data(data):
    """验证slides_data.json的schema正确性"""
    errors = []
    if not isinstance(data, dict):
        return False, ["根元素必须是dict"]
    if 'slides' not in data:
        return False, ["缺少必需的'slides'字段"]
    slides = data.get('slides', [])
    for i, slide in enumerate(slides):
        if not isinstance(slide, dict):
            errors.append(f"slide[{i}] 必须是dict")
            continue
        for field in SLIDES_DATA_SCHEMA['required_per_slide']:
            if field not in slide:
                errors.append(f"slide[{i}] 缺少必需字段'{field}'")
        slide_type = slide.get('type', '')
        if slide_type and slide_type not in SLIDES_DATA_SCHEMA['valid_types']:
            errors.append(f"slide[{i}] type='{slide_type}' 不在有效列表中")
        layout = slide.get('layout', '')
        if layout and layout not in SLIDES_DATA_SCHEMA['valid_layouts']:
            errors.append(f"slide[{i}] layout='{layout}' 不在有效列表中")
    return len(errors) == 0, errors


# ==================== 主函数 ====================

def main():
    parser = argparse.ArgumentParser(description='组会汇报 PPT 生成器 — 浙大蓝多篇版')
    parser.add_argument('--output', '-o', required=True, help='输出 .pptx 文件路径')
    parser.add_argument('--data', '-d', default='slides_data.json', help='幻灯片数据 JSON 文件')
    parser.add_argument('--layout', '-l', default=None, help='Layout选择 JSON 文件(可选)')
    parser.add_argument('--template', '-t', default=None, help='浙大蓝模板 .pptx 文件路径(可选)')
    parser.add_argument('--images-dir', '-i', default=None, help='论文图片目录(可选)')
    args = parser.parse_args()

    with open(args.data, 'r', encoding='utf-8') as f:
        data = json.load(f)

    # P2-9: Schema验证
    valid, errors = validate_slides_data(data)
    if not valid:
        print("⚠️ slides_data.json 验证失败：")
        for err in errors:
            print(f"  ❌ {err}")

    # 合并layout选择
    if args.layout:
        with open(args.layout, 'r', encoding='utf-8') as f:
            layout_data = json.load(f)
        for key, layout_id in layout_data.items():
            slide_idx = int(key.split('_')[1]) - 1
            if slide_idx < len(data.get('slides', [])):
                if 'layout' not in data['slides'][slide_idx]:
                    data['slides'][slide_idx]['layout'] = layout_id

    # P1-4 FIX: 模板克隆路线
    prs_template = None
    template_available = False
    if args.template and os.path.exists(args.template):
        try:
            prs_template = Presentation(args.template)
            template_available = True
        except Exception as e:
            print(f"[WARN] Template load failed ({e}), falling back to Route B")

    slides_data = data.get('slides', [])
    meta = data.get('meta', {})

    # ===== Route A: 从模板克隆 =====
    if template_available and prs_template is not None:
        prs = build_route_a(args.template, slides_data, meta, args.images_dir)
        prs.save(args.output)
        sys.stdout.reconfigure(encoding='utf-8')
        print(f"[OK] PPT generated: {args.output}")
        print(f"[OK] Total slides: {len(slides_data)}")
        print(f"[OK] Route: A(template)")
        return args.output

    # ===== Route B: 纯绘制（无模板时fallback） =====
    prs = Presentation()
    prs.slide_width = pt_to_emu(SLIDE_WIDTH_PT)
    prs.slide_height = pt_to_emu(SLIDE_HEIGHT_PT)

    for idx, s in enumerate(slides_data):
        page_type = s.get('type', 'content')
        page_num = idx + 1
        total_pages = len(slides_data)
        layout_id = s.get('layout', 'C4-general')
        coords = get_layout_coords(layout_id)

        # P1-5 FIX: 解析image_path
        image_path = s.get('image_path', '')
        if image_path and args.images_dir and not os.path.isabs(image_path):
            image_path = os.path.join(args.images_dir, image_path)

        if page_type == 'cover':
            create_cover_slide(prs, title=s.get('title',''), subtitle=s.get('subtitle',''),
                               presenter=meta.get('presenter', s.get('presenter','')),
                               date_str=meta.get('date', s.get('date_str','')), coords=coords)
        elif page_type == 'toc':
            create_toc_slide(prs, items=s.get('items', s.get('bullets',[])), coords=coords)
        elif page_type == 'section_separator':
            create_section_separator_slide(prs, section_num=s.get('section_num',''),
                section_label=s.get('section_label',''), items=s.get('items', s.get('bullets',[])), coords=coords)
        elif page_type == 'ending':
            create_ending_slide(prs, coords=coords)
        elif page_type == 'critique':
            create_critique_slide(prs, critiques=s.get('critiques',[]), fix_suggestions=s.get('fix_suggestions',''),
                section_num=s.get('section_num',''), section_title=s.get('section_title', s.get('title','方法论批判与局限')),
                page_num=page_num, total=total_pages, coords=coords)
        elif page_type == 'summary':
            create_summary_slide(prs, takeaway=s.get('takeaway',''), inspirations=s.get('inspirations',[]),
                next_steps=s.get('next_steps',[]), section_num=s.get('section_num',''),
                section_title=s.get('section_title', s.get('title','总结与下一步')),
                page_num=page_num, total=total_pages, coords=coords)
        elif layout_id == 'C1-info':
            create_info_slide(prs, section_num=s.get('section_num',''), section_title=s.get('title',''),
                image_path=image_path, table_data=s.get('table_data',''), keywords=s.get('keywords',[]),
                bottom_text=s.get('bottom_text', s.get('content_text','')),
                page_num=page_num, total=total_pages, coords=coords)
        elif layout_id == 'C2-image-text':
            create_image_text_slide(prs, section_num=s.get('section_num',''), section_title=s.get('title',''),
                image_path=image_path, right_col1_title=s.get('right_col1_title',''),
                right_col1_text=s.get('right_col1_text',''), right_col2_title=s.get('right_col2_title',''),
                right_col2_text=s.get('right_col2_text',''), caption=s.get('caption',''),
                page_num=page_num, total=total_pages, coords=coords)
        elif layout_id == 'C3-flowchart':
            create_flowchart_slide(prs, section_num=s.get('section_num',''), section_title=s.get('title',''),
                flow_items=s.get('flow_items', s.get('bullets',[])),
                page_num=page_num, total=total_pages, coords=coords)
        else:
            content = s.get('content_text', '')
            if not content and s.get('bullets'):
                content = '\n'.join(f"▸  {b}" for b in s['bullets'])
            create_general_content_slide(prs, section_num=s.get('section_num',''),
                section_title=s.get('title',''), content_text=content, image_path=image_path,
                page_num=page_num, total=total_pages, coords=coords)

    prs.save(args.output)
    sys.stdout.reconfigure(encoding='utf-8')
    print(f"[OK] PPT generated: {args.output}")
    print(f"[OK] Total slides: {len(slides_data)}")
    print(f"[OK] Route: B(draw)")
    return args.output


if __name__ == '__main__':
    main()
